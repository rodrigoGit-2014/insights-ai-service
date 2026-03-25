import logging
import os
import uuid
from datetime import datetime, timezone
from sqlalchemy.orm import Session
from app.core.config import settings
from app.models.report_job import ReportJob, ReportStatus
from app.services.insight_service import InsightService

logger = logging.getLogger(__name__)

class ReportService:
    def __init__(self, db: Session):
        self.db = db

    def create_job(self, company_id, format: str, report_type: str) -> ReportJob:
        job = ReportJob(
            company_id=company_id,
            format=format,
            report_type=report_type,
            status=ReportStatus.PENDING,
        )
        self.db.add(job)
        self.db.commit()
        self.db.refresh(job)
        return job

    def get_job(self, job_id) -> ReportJob:
        return self.db.query(ReportJob).filter(ReportJob.id == job_id).first()

    async def process_report(self, job_id, auth_token, start_date=None, end_date=None, department_id=None, section_id=None):
        from app.db.session import SessionLocal
        db = SessionLocal()
        try:
            job = db.query(ReportJob).filter(ReportJob.id == job_id).first()
            if not job:
                return
            job.status = ReportStatus.GENERATING
            db.commit()

            insight_svc = InsightService()
            insight = await insight_svc.generate_insight(
                auth_token=auth_token,
                company_id=str(job.company_id),
                insight_type=job.report_type,
                start_date=start_date,
                end_date=end_date,
                department_id=department_id,
                section_id=section_id,
            )

            os.makedirs(settings.REPORT_OUTPUT_DIR, exist_ok=True)
            filename = f"{job.report_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{job.format}"
            filepath = os.path.join(settings.REPORT_OUTPUT_DIR, filename)

            if job.format == "pdf":
                self._generate_pdf(filepath, insight)
            elif job.format == "docx":
                self._generate_docx(filepath, insight)
            elif job.format == "pptx":
                self._generate_pptx(filepath, insight)

            job.file_path = filepath
            job.status = ReportStatus.COMPLETED
            job.completed_at = datetime.now(timezone.utc)
            db.commit()
            logger.info(f"Report {job_id} completed: {filepath}")

        except Exception as e:
            logger.error(f"Report generation failed: {e}", exc_info=True)
            job = db.query(ReportJob).filter(ReportJob.id == job_id).first()
            if job:
                job.status = ReportStatus.FAILED
                job.error_message = str(e)
                db.commit()
        finally:
            db.close()

    def _generate_pdf(self, filepath, insight):
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib import colors
        from reportlab.lib.units import inch

        doc = SimpleDocTemplate(filepath, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=20, spaceAfter=20)
        heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=14, spaceAfter=10, textColor=colors.HexColor('#3b82f6'))
        body_style = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=10, spaceAfter=8, leading=14)

        elements = []
        data = insight.get("data", {})

        elements.append(Paragraph(data.get("titulo", "Reporte de Inteligencia Comercial"), title_style))
        elements.append(Spacer(1, 12))

        if data.get("resumen_general"):
            elements.append(Paragraph("Resumen Ejecutivo", heading_style))
            for para in data["resumen_general"].split("\n"):
                if para.strip():
                    elements.append(Paragraph(para.strip(), body_style))
            elements.append(Spacer(1, 12))

        if data.get("hallazgos_clave"):
            elements.append(Paragraph("Hallazgos Clave", heading_style))
            for h in data["hallazgos_clave"]:
                elements.append(Paragraph(f"<b>{h.get('titulo', '')}</b> [{h.get('impacto', '').upper()}]", body_style))
                elements.append(Paragraph(h.get("descripcion", ""), body_style))
                elements.append(Spacer(1, 6))

        if data.get("oportunidades_cross_selling"):
            elements.append(Paragraph("Oportunidades de Cross-Selling", heading_style))
            table_data = [["Productos", "Confianza", "Lift", "Acción"]]
            for opp in data["oportunidades_cross_selling"]:
                table_data.append([
                    " + ".join(opp.get("combinacion", [])),
                    f"{opp.get('confidence', 0)*100:.0f}%",
                    f"{opp.get('lift', 0):.1f}x",
                    opp.get("recomendacion_accion", "")[:60],
                ])
            t = Table(table_data, colWidths=[1.5*inch, 0.8*inch, 0.6*inch, 3.5*inch])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1e293b')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#334155')),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#f8fafc'), colors.white]),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 12))

        if data.get("recomendaciones_estrategicas"):
            elements.append(Paragraph("Recomendaciones Estratégicas", heading_style))
            for rec in data["recomendaciones_estrategicas"]:
                elements.append(Paragraph(f"<b>[{rec.get('prioridad', '').upper()}] {rec.get('titulo', '')}</b>", body_style))
                elements.append(Paragraph(rec.get("descripcion", ""), body_style))
                elements.append(Spacer(1, 4))

        doc.build(elements)

    def _generate_docx(self, filepath, insight):
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        data = insight.get("data", {})

        title = doc.add_heading(data.get("titulo", "Reporte de Inteligencia Comercial"), level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        if data.get("resumen_general"):
            doc.add_heading("Resumen Ejecutivo", level=1)
            doc.add_paragraph(data["resumen_general"])

        if data.get("hallazgos_clave"):
            doc.add_heading("Hallazgos Clave", level=1)
            for h in data["hallazgos_clave"]:
                p = doc.add_paragraph()
                run = p.add_run(f"{h.get('titulo', '')} ")
                run.bold = True
                badge = p.add_run(f"[{h.get('impacto', '').upper()}]")
                badge.font.size = Pt(8)
                doc.add_paragraph(h.get("descripcion", ""))

        if data.get("oportunidades_cross_selling"):
            doc.add_heading("Oportunidades de Cross-Selling", level=1)
            table = doc.add_table(rows=1, cols=4)
            table.style = 'Table Grid'
            hdr = table.rows[0].cells
            hdr[0].text = "Productos"
            hdr[1].text = "Confianza"
            hdr[2].text = "Lift"
            hdr[3].text = "Acción Recomendada"
            for opp in data["oportunidades_cross_selling"]:
                row = table.add_row().cells
                row[0].text = " + ".join(opp.get("combinacion", []))
                row[1].text = f"{opp.get('confidence', 0)*100:.0f}%"
                row[2].text = f"{opp.get('lift', 0):.1f}x"
                row[3].text = opp.get("recomendacion_accion", "")

        if data.get("recomendaciones_estrategicas"):
            doc.add_heading("Recomendaciones Estratégicas", level=1)
            for rec in data["recomendaciones_estrategicas"]:
                p = doc.add_paragraph()
                run = p.add_run(f"[{rec.get('prioridad', '').upper()}] {rec.get('titulo', '')}")
                run.bold = True
                doc.add_paragraph(rec.get("descripcion", ""))

        doc.save(filepath)

    def _generate_pptx(self, filepath, insight):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        data = insight.get("data", {})

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = data.get("titulo", "Inteligencia Comercial")
        slide.placeholders[1].text = "Generado con IA"

        # Summary slide
        if data.get("resumen_general"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Resumen Ejecutivo"
            slide.placeholders[1].text = data["resumen_general"][:800]

        # Findings slides
        for h in data.get("hallazgos_clave", [])[:5]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = h.get("titulo", "Hallazgo")
            body = slide.placeholders[1]
            body.text = h.get("descripcion", "")
            if h.get("productos_involucrados"):
                body.text += f"\n\nProductos: {', '.join(h['productos_involucrados'])}"
            body.text += f"\nImpacto: {h.get('impacto', 'N/A').upper()}"

        # Recommendations slide
        if data.get("recomendaciones_estrategicas"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Recomendaciones"
            body = slide.placeholders[1].text_frame
            body.clear()
            for rec in data["recomendaciones_estrategicas"][:6]:
                p = body.add_paragraph()
                p.text = f"[{rec.get('prioridad', '').upper()}] {rec.get('titulo', '')}: {rec.get('descripcion', '')}"
                p.font.size = Pt(14)

        prs.save(filepath)
